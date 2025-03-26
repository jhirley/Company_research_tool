import os
import json
from typing import List, Dict, Optional
from fastapi import FastAPI, Request, Form, UploadFile, File, HTTPException, Depends
from fastapi.responses import HTMLResponse, RedirectResponse, JSONResponse
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
from pydantic import BaseModel
import httpx
import re
from datetime import datetime

# Initialize FastAPI app
app = FastAPI(title="Company Research Tool")

# Mount static files
app.mount("/static", StaticFiles(directory="static"), name="static")

# Set up Jinja2 templates
templates = Jinja2Templates(directory="templates")

# Global variables to store API keys and settings
api_keys = {"openai": "", "gemini": ""}
current_companies = []
research_results = {}
prompt_categories = {}

# Load prompts from files
def load_prompts():
    global prompt_categories
    prompt_files = [f for f in os.listdir() if f.startswith("Prompts_") and f.endswith(".txt")]
    
    for file_name in prompt_files:
        # Try different encodings to handle potential encoding issues
        encodings = ['utf-8', 'latin-1', 'cp1252']
        content = None
        
        for encoding in encodings:
            try:
                with open(file_name, "r", encoding=encoding) as file:
                    content = file.read()
                    break  # If successful, break the loop
            except UnicodeDecodeError:
                continue  # Try the next encoding
        
        if content is None:
            print(f"Warning: Could not decode {file_name} with any of the attempted encodings.")
            continue
            
        category_name = file_name.replace("Prompts_", "").replace(".txt", "")
        categories = {}
        
        # Extract categories and their prompts
        current_category = None
        current_prompts = []
        
        for line in content.split("\n"):
            line = line.strip()
            if not line:
                continue
            
            # Check if this is a category header
            if line.startswith("###"):
                if current_category and current_prompts:
                    categories[current_category] = current_prompts
                current_category = line.replace("###", "").strip()
                current_prompts = []
            elif line.startswith("##"):
                if current_category and current_prompts:
                    categories[current_category] = current_prompts
                current_category = line.replace("##", "").strip()
                current_prompts = []
            elif line.startswith("#"):
                if current_category and current_prompts:
                    categories[current_category] = current_prompts
                current_category = line.replace("#", "").strip()
                current_prompts = []
            # Check if this is a numbered prompt
            elif re.match(r'^\d+\.', line) and current_category:
                # Extract the prompt number and text
                match = re.match(r'^(\d+)\.\s*(.*)', line)
                if match:
                    prompt_num = int(match.group(1))
                    prompt_text = match.group(2)
                    current_prompts.append({"number": prompt_num, "text": prompt_text})
        
        # Add the last category
        if current_category and current_prompts:
            categories[current_category] = current_prompts
            
        prompt_categories[category_name] = categories

# Save prompts to files
def save_prompts():
    for category_name, categories in prompt_categories.items():
        file_name = f"Prompts_{category_name}.txt"
        with open(file_name, "w", encoding='utf-8') as file:
            for cat_name, prompts in categories.items():
                file.write(f"### {cat_name}\n")
                for prompt in prompts:
                    file.write(f"{prompt['number']}. {prompt['text']}\n")
                file.write("\n")

# Initialize by loading prompts
@app.on_event("startup")
async def startup_event():
    load_prompts()

# Home page
@app.get("/", response_class=HTMLResponse)
async def home(request: Request):
    return templates.TemplateResponse(
        "index.html", 
        {
            "request": request,
            "api_keys": api_keys,
            "companies": current_companies
        }
    )

# API key setup
@app.post("/setup")
async def setup(
    request: Request,
    openai_key: str = Form(None),
    gemini_key: str = Form(None),
    companies: str = Form(...)
):
    # Update API keys if provided
    if openai_key:
        api_keys["openai"] = openai_key
    if gemini_key:
        api_keys["gemini"] = gemini_key
    
    # Update companies list
    global current_companies
    current_companies = [company.strip() for company in companies.split(",") if company.strip()]
    
    # Redirect to research page
    return RedirectResponse(url="/research", status_code=303)

# Research page
@app.get("/research", response_class=HTMLResponse)
async def research_page(request: Request):
    if not current_companies:
        return RedirectResponse(url="/", status_code=303)
    
    return templates.TemplateResponse(
        "research.html", 
        {
            "request": request,
            "companies": current_companies,
            "prompt_categories": prompt_categories
        }
    )

# Prompts management page
@app.get("/prompts", response_class=HTMLResponse)
async def prompts_page(request: Request):
    return templates.TemplateResponse(
        "prompts.html", 
        {
            "request": request,
            "prompt_categories": prompt_categories
        }
    )

# Update prompts
@app.post("/update_prompts")
async def update_prompts(request: Request):
    form_data = await request.form()
    
    # Process form data to update prompts
    updated_categories = {}
    
    for key, value in form_data.items():
        if key.startswith("prompt_"):
            # Format: prompt_category_name_subcategory_number
            parts = key.split("_")
            if len(parts) >= 4:
                category = parts[1]
                subcategory = "_".join(parts[2:-1])
                number = int(parts[-1])
                
                if category not in updated_categories:
                    updated_categories[category] = {}
                
                if subcategory not in updated_categories[category]:
                    updated_categories[category][subcategory] = []
                
                updated_categories[category][subcategory].append({
                    "number": number,
                    "text": value
                })
    
    # Update global prompt categories
    global prompt_categories
    prompt_categories = updated_categories
    
    # Save updated prompts to files
    save_prompts()
    
    return RedirectResponse(url="/prompts", status_code=303)

# Conduct research
@app.post("/conduct_research")
async def conduct_research(
    request: Request,
    company: str = Form(...),
    prompt_category: str = Form(...),
    subcategory: str = Form(...),
    prompt_number: int = Form(...)
):
    # Validate inputs
    if not api_keys["openai"] and not api_keys["gemini"]:
        return JSONResponse(
            status_code=400,
            content={"error": "No API key provided. Please add an OpenAI or Gemini API key on the home page."}
        )
    
    if company not in current_companies:
        return JSONResponse(
            status_code=400,
            content={"error": f"Invalid company '{company}'. Please select a valid company from the list."}
        )
    
    if prompt_category not in prompt_categories:
        return JSONResponse(
            status_code=400,
            content={"error": f"Invalid prompt category '{prompt_category}'. Please select a valid category."}
        )
    
    if subcategory not in prompt_categories[prompt_category]:
        return JSONResponse(
            status_code=400,
            content={"error": f"Invalid subcategory '{subcategory}'. Please select a valid subcategory."}
        )
    
    # Find the prompt
    prompt = None
    for p in prompt_categories[prompt_category][subcategory]:
        if p["number"] == prompt_number:
            prompt = p["text"]
            break
    
    if not prompt:
        return JSONResponse(
            status_code=400,
            content={"error": f"Invalid prompt number {prompt_number}. Please select a valid prompt."}
        )
    
    # Replace placeholder with company name
    prompt = prompt.replace("{Customer}", company)
    prompt = prompt.replace("{Customer Name}", company)
    
    # Conduct research using the appropriate API
    result = ""
    sources = []
    
    try:
        print(f"Starting research for {company} using {'OpenAI' if api_keys['openai'] else 'Gemini'} API")
        if api_keys["openai"]:
            result, sources = await conduct_research_openai(prompt, company)
        elif api_keys["gemini"]:
            result, sources = await conduct_research_gemini(prompt, company)
        print(f"Research completed successfully for {company}")
    except Exception as e:
        error_message = str(e)
        print(f"Research error: {error_message}")
        return JSONResponse(
            status_code=500,
            content={
                "error": "An error occurred while conducting research", 
                "details": error_message
            }
        )
    
    # Store research results
    if company not in research_results:
        research_results[company] = []
    
    research_results[company].append({
        "prompt": prompt,
        "result": result,
        "sources": sources,
        "category": prompt_category,
        "subcategory": subcategory
    })
    
    # Return the research results
    return {
        "company": company,
        "prompt": prompt,
        "result": result,
        "sources": sources
    }

# OpenAI research function
async def conduct_research_openai(prompt: str, company: str):
    from openai import AsyncOpenAI
    import json
    import logging
    import traceback
    
    try:
        print(f"Starting OpenAI research for company: {company}")
        print(f"Using API key: {api_keys['openai'][:5]}...{api_keys['openai'][-4:] if len(api_keys['openai']) > 10 else ''}")
        
        # Create a client instance with the API key
        client = AsyncOpenAI(api_key=api_keys["openai"])
        
        # Create a system message that emphasizes the need to avoid hallucinations
        system_message = f"""
        You are a professional business analyst researching {company}. 
        Your task is to provide accurate, factual information based on verifiable data.
        
        IMPORTANT GUIDELINES:
        1. Only provide information that you can verify from reliable sources.
        2. Clearly cite your sources for each piece of information.
        3. If you don't have enough information to answer a question, state this clearly rather than making assumptions.
        4. Avoid making predictions or speculations unless explicitly asked to do so.
        5. Format your response in a clear, professional manner.
        6. Include a "SOURCES" section at the end with numbered references to all sources used.
        """
        
        print(f"Sending prompt to OpenAI: {prompt[:50]}...")
        
        # Make the API call with the new client format
        try:
            response = await client.chat.completions.create(
                model="gpt-4",  # Using GPT-4 model (fallback from gpt-4o which might not be available)
                messages=[
                    {"role": "system", "content": system_message},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.3,  # Lower temperature for more factual responses
                max_tokens=2000
            )
            print("Successfully received response from OpenAI")
        except Exception as api_error:
            print(f"Error during OpenAI API call: {str(api_error)}")
            print(f"Error details: {traceback.format_exc()}")
            raise Exception(f"OpenAI API call failed: {str(api_error)}")
        
        # Extract the response text
        result = response.choices[0].message.content
        print(f"Received result of length: {len(result)} characters")
        
        # Extract sources from the response
        sources = []
        if "SOURCES" in result:
            try:
                sources_section = result.split("SOURCES")[1]
                source_lines = sources_section.strip().split("\n")
                for line in source_lines:
                    if line.strip():
                        # Remove any numbering or bullet points at the beginning
                        cleaned_line = re.sub(r'^[\d\.\-\*\s]+', '', line.strip())
                        sources.append(cleaned_line)
                print(f"Extracted {len(sources)} sources from the response")
            except Exception as source_error:
                print(f"Error extracting sources: {str(source_error)}")
        else:
            print("No SOURCES section found in the response")
        
        return result, sources
    except Exception as e:
        print(f"OpenAI API error: {str(e)}")
        print(f"Full error details: {traceback.format_exc()}")
        raise Exception(f"Error with OpenAI API: {str(e)}")

# Gemini research function
async def conduct_research_gemini(prompt: str, company: str):
    from google import genai
    import asyncio
    import traceback
    
    try:
        print(f"Starting Gemini research for company: {company}")
        print(f"Using API key: {api_keys['gemini'][:5]}...{api_keys['gemini'][-4:] if len(api_keys['gemini']) > 10 else ''}")
        
        # Initialize the client with the API key
        try:
            client = genai.Client(api_key=api_keys["gemini"])
            print("Successfully initialized Gemini client")
        except Exception as client_error:
            print(f"Error initializing Gemini client: {str(client_error)}")
            print(f"Error details: {traceback.format_exc()}")
            raise Exception(f"Failed to initialize Gemini client: {str(client_error)}")
        
        # Create a system message that emphasizes the need to avoid hallucinations
        system_prompt = f"""
        You are a professional business analyst researching {company}. 
        Your task is to provide accurate, factual information based on verifiable data.
        
        IMPORTANT GUIDELINES:
        1. Only provide information that you can verify from reliable sources.
        2. Clearly cite your sources for each piece of information.
        3. If you don't have enough information to answer a question, state this clearly rather than making assumptions.
        4. Avoid making predictions or speculations unless explicitly asked to do so.
        5. Format your response in a clear, professional manner.
        6. Include a "SOURCES" section at the end with numbered references to all sources used.
        """
        
        print(f"Sending prompt to Gemini: {prompt[:50]}...")
        
        # Prepare the contents with system prompt and user prompt
        contents = [
            {"role": "system", "parts": [system_prompt]},
            {"role": "user", "parts": [prompt]}
        ]
        
        # Make the API call - wrap in asyncio.to_thread since Gemini's API is synchronous
        async def generate():
            try:
                # Try with different model names to find one that works
                models_to_try = [
                    "models/gemini-1.5-pro",  # Full path format
                    "models/gemini-pro",      # Full path format
                    "gemini-1.5-pro",         # Short name format
                    "gemini-pro"              # Short name format
                ]
                
                last_error = None
                for model_name in models_to_try:
                    try:
                        print(f"Attempting to use Gemini model: {model_name}")
                        response = client.models.generate_content(
                            model=model_name,
                            contents=contents,
                            generation_config={"temperature": 0.3}  # Lower temperature for more factual responses
                        )
                        print(f"Successfully received response from Gemini model: {model_name}")
                        return response
                    except Exception as model_error:
                        last_error = model_error
                        print(f"Error with model {model_name}: {str(model_error)}")
                
                # If we get here, all models failed
                raise Exception(f"All Gemini models failed. Last error: {str(last_error)}")
            except Exception as api_error:
                print(f"Error during Gemini API call: {str(api_error)}")
                print(f"Error details: {traceback.format_exc()}")
                raise Exception(f"Gemini API call failed: {str(api_error)}")
        
        response = await asyncio.to_thread(generate)
        
        # Extract the response text
        try:
            result = response.text
            print(f"Received result of length: {len(result)} characters")
        except Exception as text_error:
            print(f"Error extracting text from Gemini response: {str(text_error)}")
            print(f"Response object: {response}")
            raise Exception(f"Failed to extract text from Gemini response: {str(text_error)}")
        
        # Extract sources from the response
        sources = []
        if "SOURCES" in result:
            try:
                sources_section = result.split("SOURCES")[1]
                source_lines = sources_section.strip().split("\n")
                for line in source_lines:
                    if line.strip():
                        # Remove any numbering or bullet points at the beginning
                        cleaned_line = re.sub(r'^[\d\.\-\*\s]+', '', line.strip())
                        sources.append(cleaned_line)
                print(f"Extracted {len(sources)} sources from the response")
            except Exception as source_error:
                print(f"Error extracting sources: {str(source_error)}")
        else:
            print("No SOURCES section found in the response")
        
        return result, sources
    except Exception as e:
        print(f"Gemini API error: {str(e)}")
        print(f"Full error details: {traceback.format_exc()}")
        raise Exception(f"Error with Gemini API: {str(e)}")

# Export results
@app.get("/export/{company}")
async def export_results(company: str):
    if company not in research_results:
        raise HTTPException(status_code=404, detail="No research results found for this company")
    
    # Create a formatted export
    export_data = {
        "company": company,
        "research_results": research_results[company]
    }
    
    return JSONResponse(content=export_data)

# Export results in different formats
@app.get("/export/{company}/{format}")
async def export_results_format(company: str, format: str):
    if company not in research_results:
        raise HTTPException(status_code=404, detail="No research results found for this company")
    
    # Create a formatted export
    export_data = {
        "company": company,
        "research_results": research_results[company]
    }
    
    if format == "json":
        return JSONResponse(content=export_data)
    elif format == "word":
        from docx import Document
        from io import BytesIO
        from fastapi.responses import StreamingResponse
        
        # Create a new Word document
        doc = Document()
        doc.add_heading(f"Research Results for {company}", 0)
        
        # Add each research result to the document
        for result in research_results[company]:
            doc.add_heading(f"Category: {result['category']} - {result['subcategory']}", 1)
            doc.add_paragraph(f"Prompt: {result['prompt']}")
            doc.add_paragraph(result['result'])
            if result.get('sources'):
                doc.add_heading("Sources:", 2)
                for i, source in enumerate(result['sources']):
                    doc.add_paragraph(f"{i+1}. {source}")
            doc.add_page_break()
        
        # Save the document to a BytesIO object
        f = BytesIO()
        doc.save(f)
        f.seek(0)
        
        # Return the document as a downloadable file
        return StreamingResponse(
            f, 
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers={"Content-Disposition": f"attachment; filename={company}_research.docx"}
        )
    elif format == "pdf":
        from fpdf import FPDF
        from io import BytesIO
        from fastapi.responses import StreamingResponse
        
        # Create a new PDF document
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()
        
        # Set fonts
        pdf.set_font("Arial", "B", 16)
        pdf.cell(0, 10, f"Research Results for {company}", 0, 1, "C")
        
        # Add each research result to the document
        for result in research_results[company]:
            pdf.add_page()
            pdf.set_font("Arial", "B", 14)
            pdf.cell(0, 10, f"Category: {result['category']} - {result['subcategory']}", 0, 1)
            
            pdf.set_font("Arial", "B", 12)
            pdf.cell(0, 10, "Prompt:", 0, 1)
            
            pdf.set_font("Arial", "", 12)
            pdf.multi_cell(0, 10, result['prompt'])
            
            pdf.set_font("Arial", "B", 12)
            pdf.cell(0, 10, "Result:", 0, 1)
            
            pdf.set_font("Arial", "", 12)
            pdf.multi_cell(0, 10, result['result'])
            
            if result.get('sources'):
                pdf.set_font("Arial", "B", 12)
                pdf.cell(0, 10, "Sources:", 0, 1)
                
                pdf.set_font("Arial", "", 12)
                for i, source in enumerate(result['sources']):
                    pdf.multi_cell(0, 10, f"{i+1}. {source}")
        
        # Save the PDF to a BytesIO object
        f = BytesIO()
        f.write(pdf.output(dest='S').encode('latin1'))
        f.seek(0)
        
        # Return the PDF as a downloadable file
        return StreamingResponse(
            f, 
            media_type="application/pdf",
            headers={"Content-Disposition": f"attachment; filename={company}_research.pdf"}
        )
    elif format == "excel":
        import pandas as pd
        from io import BytesIO
        from fastapi.responses import StreamingResponse
        
        # Create a new Excel workbook
        f = BytesIO()
        writer = pd.ExcelWriter(f, engine='xlsxwriter')
        
        # Create a summary sheet
        summary_data = []
        for result in research_results[company]:
            summary_data.append({
                "Category": result['category'],
                "Subcategory": result['subcategory'],
                "Prompt": result['prompt'][:100] + "..." if len(result['prompt']) > 100 else result['prompt']
            })
        
        summary_df = pd.DataFrame(summary_data)
        summary_df.to_excel(writer, sheet_name="Summary", index=False)
        
        # Create a sheet for each research result
        for i, result in enumerate(research_results[company]):
            sheet_name = f"{result['category']}_{i+1}"
            if len(sheet_name) > 31:  # Excel sheet name length limit
                sheet_name = sheet_name[:31]
            
            # Create a DataFrame for this result
            result_data = {
                "Field": ["Category", "Subcategory", "Prompt", "Result"],
                "Value": [
                    result['category'],
                    result['subcategory'],
                    result['prompt'],
                    result['result']
                ]
            }
            
            # Add sources if available
            if result.get('sources'):
                for i, source in enumerate(result['sources']):
                    result_data["Field"].append(f"Source {i+1}")
                    result_data["Value"].append(source)
            
            result_df = pd.DataFrame(result_data)
            result_df.to_excel(writer, sheet_name=sheet_name, index=False)
        
        # Save the workbook
        writer.close()
        f.seek(0)
        
        # Return the Excel file as a downloadable file
        return StreamingResponse(
            f, 
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            headers={"Content-Disposition": f"attachment; filename={company}_research.xlsx"}
        )
    elif format == "pptx":
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from io import BytesIO
        from fastapi.responses import StreamingResponse
        from collections import defaultdict
        
        # Create a new PowerPoint presentation
        prs = Presentation()
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = f"Research Results for {company}"
        subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d')}"
        
        # Organize research results by subcategory
        subcategory_results = defaultdict(list)
        for result in research_results[company]:
            key = f"{result['category']} - {result['subcategory']}"
            subcategory_results[key].append(result)
        
        # Add content slides - one per subcategory
        content_slide_layout = prs.slide_layouts[1]  # Title and content layout
        
        for subcategory, results in subcategory_results.items():
            slide = prs.slides.add_slide(content_slide_layout)
            title = slide.shapes.title
            title.text = subcategory
            
            # Add content placeholder
            content = slide.placeholders[1]
            tf = content.text_frame
            
            # Add each prompt and result for this subcategory
            for result in results:
                p = tf.add_paragraph()
                p.text = f"Prompt: {result['prompt'][:100]}..." if len(result['prompt']) > 100 else f"Prompt: {result['prompt']}"
                p.font.bold = True
                
                p = tf.add_paragraph()
                p.text = result['result'][:500] + "..." if len(result['result']) > 500 else result['result']
                
                # Add sources if available
                if result.get('sources') and len(result['sources']) > 0:
                    p = tf.add_paragraph()
                    p.text = "Sources:"
                    p.font.bold = True
                    
                    for source in result['sources'][:3]:  # Limit to first 3 sources to save space
                        p = tf.add_paragraph()
                        p.text = source[:100] + "..." if len(source) > 100 else source
                        p.level = 1  # Indent sources
                
                # Add a separator between results (except for the last one)
                if result != results[-1]:
                    p = tf.add_paragraph()
                    p.text = "---"
        
        # Save the presentation to a BytesIO object
        f = BytesIO()
        prs.save(f)
        f.seek(0)
        
        # Return the PowerPoint file as a downloadable file
        return StreamingResponse(
            f, 
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f"attachment; filename={company}_research.pptx"}
        )
    else:
        raise HTTPException(status_code=400, detail=f"Unsupported export format: {format}")

# Research all prompts for a company
@app.post("/conduct_research_all")
async def conduct_research_all(
    request: Request,
    company: str = Form(...)
):
    # Validate inputs
    if not api_keys["openai"] and not api_keys["gemini"]:
        return JSONResponse(
            status_code=400,
            content={"error": "No API key provided. Please add an OpenAI or Gemini API key on the home page."}
        )
    
    if company not in current_companies:
        return JSONResponse(
            status_code=400,
            content={"error": f"Invalid company '{company}'. Please select a valid company from the list."}
        )
    
    # Initialize results storage if needed
    if company not in research_results:
        research_results[company] = []
    
    # Track all results
    all_results = []
    errors = []
    
    # Iterate through all categories, subcategories, and prompts
    for category in prompt_categories:
        for subcategory in prompt_categories[category]:
            for prompt_obj in prompt_categories[category][subcategory]:
                prompt_number = prompt_obj["number"]
                prompt_text = prompt_obj["text"]
                
                # Replace placeholder with company name
                prompt = prompt_text.replace("{Customer}", company)
                prompt = prompt.replace("{Customer Name}", company)
                
                try:
                    # Conduct research using the appropriate API
                    result = ""
                    sources = []
                    
                    print(f"Starting research for {company}, {category}/{subcategory}, prompt #{prompt_number}")
                    if api_keys["openai"]:
                        result, sources = await conduct_research_openai(prompt, company)
                    elif api_keys["gemini"]:
                        result, sources = await conduct_research_gemini(prompt, company)
                    
                    # Store research results (avoid duplicates)
                    existing_entry = False
                    for entry in research_results[company]:
                        if (entry["category"] == category and 
                            entry["subcategory"] == subcategory and 
                            entry["prompt"] == prompt):
                            existing_entry = True
                            entry["result"] = result
                            entry["sources"] = sources
                            break
                    
                    if not existing_entry:
                        research_results[company].append({
                            "prompt": prompt,
                            "result": result,
                            "sources": sources,
                            "category": category,
                            "subcategory": subcategory
                        })
                    
                    # Add to results for this batch
                    all_results.append({
                        "category": category,
                        "subcategory": subcategory,
                        "prompt_number": prompt_number,
                        "prompt": prompt,
                        "result": result,
                        "sources": sources
                    })
                    
                except Exception as e:
                    error_message = str(e)
                    print(f"Research error for {category}/{subcategory}, prompt #{prompt_number}: {error_message}")
                    errors.append({
                        "category": category,
                        "subcategory": subcategory,
                        "prompt_number": prompt_number,
                        "error": error_message
                    })
    
    # Return the research results
    return {
        "company": company,
        "results": all_results,
        "errors": errors,
        "total_successful": len(all_results),
        "total_errors": len(errors)
    }

# API status endpoint
@app.get("/api/status")
async def api_status():
    """Return the status of the API and available services"""
    return {
        "status": "ok",
        "version": "1.0.0",
        "services": {
            "openai": bool(api_keys["openai"]),
            "gemini": bool(api_keys["gemini"])
        },
        "companies": len(current_companies),
        "prompt_categories": len(prompt_categories)
    }

# Test API endpoint for diagnosing issues
@app.get("/api/test")
async def test_api():
    """Test the API connections and return diagnostic information"""
    results = {
        "openai": {"status": "not_tested", "error": None},
        "gemini": {"status": "not_tested", "error": None}
    }
    
    # Test OpenAI API
    if api_keys["openai"]:
        try:
            from openai import AsyncOpenAI
            client = AsyncOpenAI(api_key=api_keys["openai"])
            # Simple test request
            response = await client.chat.completions.create(
                model="gpt-3.5-turbo",  # Using a simpler model for testing
                messages=[{"role": "user", "content": "Hello, this is a test."}],
                max_tokens=10
            )
            results["openai"] = {
                "status": "ok",
                "model": "gpt-3.5-turbo",
                "response": response.choices[0].message.content
            }
        except Exception as e:
            results["openai"] = {
                "status": "error",
                "error": str(e),
                "key_length": len(api_keys["openai"]) if api_keys["openai"] else 0
            }
    
    # Test Gemini API
    if api_keys["gemini"]:
        try:
            from google import genai
            import asyncio
            
            # Initialize client
            client = genai.Client(api_key=api_keys["gemini"])
            
            # List available models
            result["step"] = "Listing available models"
            
            async def list_models():
                try:
                    models = client.models.list()
                    return [model.name for model in models]
                except Exception as e:
                    return [f"Error listing models: {str(e)}"]
            
            result["available_models"] = await asyncio.to_thread(list_models)
            
            # Test with a simple prompt
            result["step"] = "Testing content generation with multiple model formats"
            
            # Try different model formats
            models_to_try = [
                "models/gemini-1.5-pro",  # Full path format
                "models/gemini-pro",      # Full path format
                "gemini-1.5-pro",         # Short name format
                "gemini-pro"              # Short name format
            ]
            
            result["model_tests"] = {}
            working_model = None
            
            async def test_model(model_name):
                try:
                    response = client.models.generate_content(
                        model=model_name,
                        contents="Say hello in exactly 5 words."
                    )
                    return {
                        "status": "success",
                        "response": response.text
                    }
                except Exception as e:
                    return {
                        "status": "error",
                        "error": str(e)
                    }
            
            for model in models_to_try:
                # Create a function that captures the current model name
                async def test_current_model(model_name=model):
                    return await test_model(model_name)
                
                # Run the test for this model
                model_result = await asyncio.to_thread(lambda m=model: client.models.generate_content(
                    model=m,
                    contents="Say hello in exactly 5 words."
                ))
                
                try:
                    result["model_tests"][model] = {
                        "status": "success",
                        "response": model_result.text
                    }
                    working_model = model
                    result["working_model"] = model
                    # Break once we find a working model
                    break
                except Exception as e:
                    result["model_tests"][model] = {
                        "status": "error",
                        "error": str(e)
                    }
            
            if working_model:
                result["test_response"] = result["model_tests"][working_model]["response"]
                result["status"] = "success"
            else:
                result["status"] = "error"
                result["error"] = "All model formats failed"
            
        except Exception as e:
            result["status"] = "error"
            result["error"] = str(e)
            result["key_length"] = len(api_keys["gemini"]) if api_keys["gemini"] else 0
    
    return results

# Manual test endpoint for more detailed testing
@app.get("/api/manual-test")
async def manual_test(provider: str):
    """Run a manual test for a specific provider with detailed output"""
    import traceback
    
    if provider not in ["openai", "gemini"]:
        return {"error": f"Invalid provider: {provider}. Must be 'openai' or 'gemini'"}
    
    if not api_keys[provider]:
        return {"error": f"No API key configured for {provider}"}
    
    result = {
        "provider": provider,
        "key_length": len(api_keys[provider]),
        "key_preview": f"{api_keys[provider][:5]}...{api_keys[provider][-4:]}" if len(api_keys[provider]) > 10 else "[too short]"
    }
    
    try:
        if provider == "openai":
            from openai import AsyncOpenAI
            
            # Create client
            result["step"] = "Creating AsyncOpenAI client"
            client = AsyncOpenAI(api_key=api_keys["openai"])
            
            # List available models
            result["step"] = "Listing available models"
            models = await client.models.list()
            result["available_models"] = [model.id for model in models.data]
            
            # Test with simple completion
            result["step"] = "Testing chat completion"
            response = await client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[{"role": "user", "content": "Say hello in exactly 5 words."}],
                max_tokens=20
            )
            result["test_response"] = response.choices[0].message.content
            result["status"] = "success"
            
        elif provider == "gemini":
            from google import genai
            import asyncio
            
            # Create client
            result["step"] = "Creating Gemini client"
            client = genai.Client(api_key=api_keys["gemini"])
            
            # List available models
            result["step"] = "Listing available models"
            
            async def list_models():
                try:
                    models = client.models.list()
                    return [model.name for model in models]
                except Exception as e:
                    return [f"Error listing models: {str(e)}"]
            
            result["available_models"] = await asyncio.to_thread(list_models)
            
            # Test with simple generation
            result["step"] = "Testing content generation with multiple model formats"
            
            # Try different model formats
            models_to_try = [
                "models/gemini-1.5-pro",  # Full path format
                "models/gemini-pro",      # Full path format
                "gemini-1.5-pro",         # Short name format
                "gemini-pro"              # Short name format
            ]
            
            result["model_tests"] = {}
            working_model = None
            
            async def test_model(model_name):
                try:
                    response = client.models.generate_content(
                        model=model_name,
                        contents="Say hello in exactly 5 words."
                    )
                    return {
                        "status": "success",
                        "response": response.text
                    }
                except Exception as e:
                    return {
                        "status": "error",
                        "error": str(e)
                    }
            
            for model in models_to_try:
                # Create a function that captures the current model name
                async def test_current_model(model_name=model):
                    return await test_model(model_name)
                
                # Run the test for this model
                model_result = await asyncio.to_thread(lambda m=model: client.models.generate_content(
                    model=m,
                    contents="Say hello in exactly 5 words."
                ))
                
                try:
                    result["model_tests"][model] = {
                        "status": "success",
                        "response": model_result.text
                    }
                    working_model = model
                    result["working_model"] = model
                    # Break once we find a working model
                    break
                except Exception as e:
                    result["model_tests"][model] = {
                        "status": "error",
                        "error": str(e)
                    }
            
            if working_model:
                result["test_response"] = result["model_tests"][working_model]["response"]
                result["status"] = "success"
            else:
                result["status"] = "error"
                result["error"] = "All model formats failed"
            
    except Exception as e:
        result["status"] = "error"
        result["error"] = str(e)
        result["traceback"] = traceback.format_exc()
    
    return result

# Diagnostics page
@app.get("/diagnostics", response_class=HTMLResponse)
async def diagnostics_page(request: Request):
    return templates.TemplateResponse(
        "diagnostics.html",
        {"request": request}
    )

# Run the application
if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
