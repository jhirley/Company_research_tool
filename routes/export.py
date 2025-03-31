from fastapi import APIRouter, HTTPException
from fastapi.responses import JSONResponse, StreamingResponse
from io import BytesIO
from datetime import datetime
from collections import defaultdict
from models.models import research_results
import logging
import json
import os

# Import docx related modules at the top level
from docx import Document
from docx.shared import RGBColor, Pt, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH

# Import FPDF for PDF generation
from fpdf import FPDF

# Import PowerPoint modules
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Create a backup file to store research results between requests
RESEARCH_BACKUP_FILE = "research_results_backup.json"

# Function to load research results from backup file
def load_research_results():
    global research_results
    try:
        if os.path.exists(RESEARCH_BACKUP_FILE):
            with open(RESEARCH_BACKUP_FILE, 'r') as f:
                loaded_results = json.load(f)
                # Update the global research_results dictionary
                research_results.update(loaded_results)
                logger.info(f"Loaded research results from backup file. Companies: {list(research_results.keys())}")
    except Exception as e:
        logger.error(f"Error loading research results from backup: {e}")

# Function to save research results to backup file
def save_research_results():
    try:
        with open(RESEARCH_BACKUP_FILE, 'w') as f:
            json.dump(research_results, f)
        logger.info(f"Saved research results to backup file. Companies: {list(research_results.keys())}")
    except Exception as e:
        logger.error(f"Error saving research results to backup: {e}")

# Load research results when module is imported
load_research_results()

router = APIRouter()

@router.get("/export/{company}")
async def export_results(company: str):
    logger.info(f"Export requested for company: {company}")
    logger.info(f"Available companies in research_results: {list(research_results.keys())}")
    
    if company not in research_results:
        logger.error(f"No research results found for company: {company}")
        raise HTTPException(status_code=404, detail="No research results found for this company")
    
    # Create a formatted export
    export_data = {
        "company": company,
        "research_results": research_results[company]
    }
    
    logger.info(f"Successfully exported JSON data for company: {company}")
    return JSONResponse(content=export_data)

@router.get("/export/{company}/{format}")
async def export_results_format(company: str, format: str):
    logger.info(f"Export requested for company: {company} in format: {format}")
    logger.info(f"Available companies in research_results: {list(research_results.keys())}")

    if company not in research_results:
        raise HTTPException(status_code=404, detail="No research results found for this company")
    
    # Create a formatted export
    export_data = {
        "company": company,
        "research_results": research_results[company]
    }
    
    if format == "json":
        return JSONResponse(content=export_data)
    elif format == "word":
        # Create a new Word document with simple formatting
        from io import BytesIO
        doc = Document()
        
        # Add title
        title = doc.add_heading(f"Research Results for {company}", 0)
        
        # Add each research result to the document with simple styling
        for result in research_results[company]:
            # Add category heading
            doc.add_heading(f"Category: {result['promptCategory']} - {result['subcategory']}", 1)
            
            # Add prompt
            prompt_para = doc.add_paragraph()
            prompt_run = prompt_para.add_run("Prompt: ")
            prompt_run.bold = True
            prompt_para.add_run(result['prompt'])
            
            # Add result
            result_para = doc.add_paragraph()
            result_run = result_para.add_run("Result: ")
            result_run.bold = True
            result_para.add_run(result['result'])
            
            # Add sources if available
            if result.get('sources'):
                doc.add_heading("Sources:", 2)
                
                for i, source in enumerate(result['sources']):
                    doc.add_paragraph(f"{i+1}. {source}")
            
            # Add page break between results
            doc.add_page_break()
        
        # Save the document to a BytesIO object
        word_buffer = BytesIO()
        doc.save(word_buffer)
        word_buffer.seek(0)
        
        # Return the document as a downloadable file
        return StreamingResponse(
            word_buffer, 
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers={"Content-Disposition": f"attachment; filename={company}_research.docx"}
        )
    elif format == "pdf":
        # Use ReportLab for better Unicode support
        from io import BytesIO
        from reportlab.lib.pagesizes import letter
        from reportlab.lib import colors
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.enums import TA_CENTER, TA_LEFT
        from reportlab.lib.units import inch
        from reportlab.platypus.flowables import Flowable
        
        class GradientBackground(Flowable):
            """Custom flowable for a gradient background"""
            def __init__(self, width, height, start_color, end_color):
                Flowable.__init__(self)
                self.width = width
                self.height = height
                self.start_color = start_color
                self.end_color = end_color
            
            def draw(self):
                # Draw gradient rectangles
                canvas = self.canv
                steps = 20
                for i in range(steps):
                    # Calculate color for this step
                    r = self.start_color.red + (self.end_color.red - self.start_color.red) * (i / steps)
                    g = self.start_color.green + (self.end_color.green - self.start_color.green) * (i / steps)
                    b = self.start_color.blue + (self.end_color.blue - self.start_color.blue) * (i / steps)
                    
                    color = colors.Color(r, g, b)
                    canvas.setFillColor(color)
                    y_pos = self.height - (i * (self.height/steps))
                    canvas.rect(0, y_pos - self.height/steps, self.width, self.height/steps, fill=1, stroke=0)
        
        # Define theme colors
        THEME_PRIMARY = colors.Color(243/255, 88/255, 88/255)  # #f35858 - accent color
        THEME_GRADIENT_START = colors.Color(255/255, 112/255, 136/255)  # #FF7088 - gradient start
        THEME_GRADIENT_END = colors.Color(242/255, 182/255, 157/255)  # #F2B69D - gradient end
        THEME_BACKGROUND = colors.Color(245/255, 247/255, 250/255)  # #f5f7fa - background color
        THEME_TEXT = colors.Color(82/255, 87/255, 92/255)  # #52575c - text color
        
        # Create a BytesIO buffer for the PDF
        pdf_buffer = BytesIO()
        
        try:
            # Create the PDF document
            doc = SimpleDocTemplate(
                pdf_buffer,
                pagesize=letter,
                rightMargin=0.5*inch,
                leftMargin=0.5*inch,
                topMargin=0.5*inch,
                bottomMargin=0.5*inch
            )
            
            # Create styles
            styles = getSampleStyleSheet()
            
            # Define custom styles - check if they exist first
            custom_styles = {
                'CustomTitle': ParagraphStyle(
                    name='CustomTitle',
                    fontName='Helvetica-Bold',
                    fontSize=16,
                    alignment=TA_CENTER,
                    textColor=colors.white,
                    spaceAfter=12
                ),
                'Heading1': ParagraphStyle(
                    name='Heading1',
                    fontName='Helvetica-Bold',
                    fontSize=14,
                    textColor=THEME_TEXT,
                    spaceAfter=6
                ),
                'Heading2': ParagraphStyle(
                    name='Heading2',
                    fontName='Helvetica-Bold',
                    fontSize=12,
                    textColor=THEME_TEXT,
                    spaceAfter=6
                ),
                'Normal': ParagraphStyle(
                    name='Normal',
                    fontName='Helvetica',
                    fontSize=10,
                    leading=12,
                    spaceAfter=6
                ),
                'SourceItem': ParagraphStyle(
                    name='SourceItem',
                    fontName='Helvetica',
                    fontSize=8,
                    textColor=colors.gray
                )
            }
            
            # Add custom styles if they don't exist
            for style_name, style in custom_styles.items():
                if style_name not in styles:
                    styles.add(style)
            
            # Create content elements
            content = []
            
            # Add title with gradient background
            gradient = GradientBackground(letter[0], 1*inch, THEME_GRADIENT_START, THEME_GRADIENT_END)
            content.append(gradient)
            
            # Add title
            title = Paragraph(f"<font color='white'>Research Results for {company}</font>", styles['CustomTitle'])
            content.append(title)
            content.append(Spacer(1, 0.5*inch))
            
            # Process each research result
            for result in research_results[company]:
                # Add category heading
                category = Paragraph(f"Category: {result['promptCategory']} - {result['subcategory']}", styles['Heading1'])
                content.append(category)
                
                # Add prompt
                prompt_heading = Paragraph("Prompt:", styles['Heading2'])
                content.append(prompt_heading)
                
                # Clean and escape any XML/HTML special characters in the text
                def clean_text(text):
                    if not text:
                        return ""
                    # Convert to string if not already
                    if not isinstance(text, str):
                        text = str(text)
                    # Replace XML/HTML special characters
                    text = text.replace('&', '&amp;')
                    text = text.replace('<', '&lt;')
                    text = text.replace('>', '&gt;')
                    text = text.replace('"', '&quot;')
                    text = text.replace("'", '&#39;')
                    return text
                
                prompt_text = Paragraph(clean_text(result['prompt']), styles['Normal'])
                content.append(prompt_text)
                
                # Add result
                result_heading = Paragraph("Result:", styles['Heading2'])
                content.append(result_heading)
                
                result_text = Paragraph(clean_text(result['result']), styles['Normal'])
                content.append(result_text)
                
                # Add sources if available
                if result.get('sources'):
                    sources_heading = Paragraph("Sources:", styles['Heading2'])
                    content.append(sources_heading)
                    
                    for i, source in enumerate(result['sources']):
                        source_text = Paragraph(f"{i+1}. {clean_text(source)}", styles['Normal'])
                        content.append(source_text)
                
                # Add page break between results
                content.append(PageBreak())
            
            # Remove the last page break if it exists
            if content and isinstance(content[-1], PageBreak):
                content.pop()
                
            # Build the PDF
            doc.build(content)
            pdf_buffer.seek(0)
            
        except Exception as e:
            # If ReportLab fails, create a simple PDF with FPDF as fallback
            print(f"ReportLab PDF generation failed: {str(e)}", flush=True)
            import traceback
            traceback.print_exc()
            
            # Reset buffer
            pdf_buffer = BytesIO()
            
            # Use FPDF as fallback with minimal content
            from fpdf import FPDF
            fallback_pdf = FPDF()
            fallback_pdf.add_page()
            fallback_pdf.set_font("Arial", "B", 16)
            fallback_pdf.cell(0, 10, "Research Results", 0, 1, "C")
            fallback_pdf.set_font("Arial", "", 12)
            fallback_pdf.cell(0, 10, f"Results for {company}", 0, 1)
            fallback_pdf.multi_cell(0, 10, "PDF generation encountered an issue with special characters. Please try exporting to DOCX or PPTX format for better results.")
            
            pdf_buffer.write(fallback_pdf.output(dest='S').encode('latin1'))
            pdf_buffer.seek(0)
        
        # Return the PDF as a downloadable file
        return StreamingResponse(
            pdf_buffer, 
            media_type="application/pdf",
            headers={"Content-Disposition": f"attachment; filename={company}_research.pdf"}
        )
    elif format == "excel":
        import xlsxwriter
        from io import BytesIO
        
        # Create a new Excel workbook directly with xlsxwriter
        excel_buffer = BytesIO()
        workbook = xlsxwriter.Workbook(excel_buffer)
        
        # Define theme colors based on the HTML theme (same as in PPTX export)
        THEME_PRIMARY = '#f35858'  # accent color
        THEME_GRADIENT_START = '#FF7088'  # gradient start
        THEME_GRADIENT_END = '#F2B69D'  # gradient end
        THEME_BACKGROUND = '#f5f7fa'  # background color
        THEME_TEXT = '#52575c'  # text color
        
        # Create a gradient fill format for headers
        header_format = workbook.add_format({
            'bold': True,
            'font_color': 'white',
            'bg_color': THEME_GRADIENT_START,
            'border': 1,
            'align': 'center',
            'valign': 'vcenter'
        })
        
        # Create a format for the title
        title_format = workbook.add_format({
            'bold': True,
            'font_size': 16,
            'font_color': 'white',
            'bg_color': THEME_GRADIENT_START,
            'border': 1,
            'align': 'center',
            'valign': 'vcenter'
        })
        
        # Create a summary sheet
        summary_worksheet = workbook.add_worksheet("Summary")
        
        # Add a title row with gradient styling
        summary_worksheet.merge_range('A1:C1', f"Research Results for {company}", title_format)
        summary_worksheet.set_row(0, 30)  # Make the title row taller
        
        # Add header row
        headers = ["Category", "Subcategory", "Prompt"]
        for col_num, header in enumerate(headers):
            summary_worksheet.write(1, col_num, header, header_format)
        
        # Add data rows
        row = 2
        for result in research_results[company]:
            prompt_text = result['prompt'][:100] + "..." if len(result['prompt']) > 100 else result['prompt']
            summary_worksheet.write(row, 0, result['promptCategory'])
            summary_worksheet.write(row, 1, result['subcategory'])
            summary_worksheet.write(row, 2, prompt_text)
            row += 1
        
        # Adjust column widths
        summary_worksheet.set_column('A:A', 20)
        summary_worksheet.set_column('B:B', 20)
        summary_worksheet.set_column('C:C', 50)
        
        # Create a sheet for each research result
        for i, result in enumerate(research_results[company]):
            sheet_name = f"{result['promptCategory']}_{i+1}"
            if len(sheet_name) > 31:  # Excel sheet name length limit
                sheet_name = sheet_name[:31]
            
            # Create a worksheet for this result
            worksheet = workbook.add_worksheet(sheet_name)
            
            # Add a title with gradient styling
            worksheet.merge_range('A1:B1', f"{result['promptCategory']} - {result['subcategory']}", title_format)
            worksheet.set_row(0, 30)  # Make the title row taller
            
            # Format the header row
            headers = ["Field", "Value"]
            for col_num, header in enumerate(headers):
                worksheet.write(1, col_num, header, header_format)
            
            # Create a format for field names
            field_format = workbook.add_format({
                'bold': True,
                'bg_color': THEME_BACKGROUND,
                'border': 1,
                'align': 'left',
                'valign': 'vcenter'
            })
            
            # Add the data rows
            fields = ["Category", "Subcategory", "Prompt", "Result"]
            values = [
                result['promptCategory'],
                result['subcategory'],
                result['prompt'],
                result['result']
            ]
            
            # Add sources if available
            if result.get('sources'):
                for j, source in enumerate(result['sources']):
                    fields.append(f"Source {j+1}")
                    values.append(source)
            
            # Write the data to the worksheet
            for row_num, (field, value) in enumerate(zip(fields, values)):
                worksheet.write(row_num + 2, 0, field, field_format)
                worksheet.write(row_num + 2, 1, value)
            
            # Adjust column widths
            worksheet.set_column('A:A', 15)
            worksheet.set_column('B:B', 70)
        
        # Save the workbook
        workbook.close()
        excel_buffer.seek(0)
        
        # Return the Excel file as a downloadable file
        return StreamingResponse(
            excel_buffer, 
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            headers={"Content-Disposition": f"attachment; filename={company}_research.xlsx"}
        )
    elif format == "pptx":
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
        from io import BytesIO
        from collections import defaultdict
        
        # Create a new PowerPoint presentation
        prs = Presentation()
        
        # Define theme colors based on the HTML theme
        THEME_PRIMARY = RGBColor(243, 88, 88)  # #f35858 - accent color
        THEME_GRADIENT_START = RGBColor(255, 112, 136)  # #FF7088 - gradient start
        THEME_GRADIENT_END = RGBColor(242, 182, 157)  # #F2B69D - gradient end
        THEME_BACKGROUND = RGBColor(245, 247, 250)  # #f5f7fa - background color
        THEME_TEXT = RGBColor(82, 87, 92)  # #52575c - text color
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        
        # Apply theme colors to title slide
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        # Set title text and formatting
        title.text = f"Research Results for {company}"
        title_tf = title.text_frame
        title_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_run = title_tf.paragraphs[0].runs[0]
        title_run.font.color.rgb = THEME_PRIMARY
        title_run.font.bold = True
        title_run.font.size = Pt(36)
        
        # Set subtitle text and formatting
        subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d')}"
        subtitle_tf = subtitle.text_frame
        subtitle_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        subtitle_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        subtitle_run = subtitle_tf.paragraphs[0].runs[0]
        subtitle_run.font.color.rgb = THEME_TEXT
        subtitle_run.font.size = Pt(20)
        
        # Apply background fill to title slide
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = THEME_BACKGROUND
        
        # Add content slides - one per prompt/result
        blank_slide_layout = prs.slide_layouts[6]  # Blank layout for more control
        
        for result in research_results[company]:
            # Create a new slide for each prompt/result
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Apply background fill to content slide
            background = slide.background
            fill = background.fill
            fill.solid()  # This must be called before setting fore_color
            fill.fore_color.rgb = THEME_BACKGROUND
            
            # Add a title text box manually
            left = Inches(0.5)
            top = Inches(0.5)
            width = Inches(9.0)
            height = Inches(1.0)
            
            title_box = slide.shapes.add_textbox(left, top, width, height)
            title_tf = title_box.text_frame
            title_tf.word_wrap = True
            title_tf.auto_size = MSO_AUTO_SIZE.NONE  # Disable auto-sizing
            
            p = title_tf.add_paragraph()
            p.clear()
            run = p.add_run()
            run.text = f"{result['promptCategory']} - {result['subcategory']}"
            
            # Apply paragraph formatting
            p.alignment = PP_ALIGN.LEFT
            
            # Apply run formatting
            font = run.font
            font.color.rgb = THEME_PRIMARY
            font.bold = True
            font.name = 'Calibri'
            font.size = Pt(28)
            
            # Add content text box with fixed position and size
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(9.0)
            height = Inches(5.0)
            
            content_box = slide.shapes.add_textbox(left, top, width, height)
            tf = content_box.text_frame
            tf.word_wrap = True
            
            # Explicitly disable auto-sizing to ensure font sizes are preserved
            tf.auto_size = MSO_AUTO_SIZE.NONE
            tf.word_wrap = True
            
            # Add prompt with styling - use the first paragraph to avoid blank line
            if hasattr(tf, 'paragraphs') and len(tf.paragraphs) > 0:
                p = tf.paragraphs[0]  # Use the first paragraph that already exists
            else:
                p = tf.add_paragraph()
            
            # Clear any existing text and add new text as a fresh run
            p.clear()
            run = p.add_run()
            run.text = f"Prompt: {result['prompt'][:100]}..." if len(result['prompt']) > 100 else f"Prompt: {result['prompt']}"
            
            # Apply paragraph formatting
            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(0)  # No space before first paragraph
            p.space_after = Pt(6)
            
            # Apply run formatting
            font = run.font
            font.bold = True
            font.size = Pt(14)
            font.name = 'Calibri'
            font.color.rgb = THEME_GRADIENT_START
            
            # Add result with styling
            p = tf.add_paragraph()
            p.clear()
            run = p.add_run()
            run.text = result['result'][:1500] + "..." if len(result['result']) > 1500 else result['result']
            
            # Apply paragraph formatting
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(4)
            
            # Apply run formatting
            font = run.font
            font.size = Pt(11)  # Set to 11pt
            font.name = 'Calibri'
            font.color.rgb = THEME_TEXT
            
            # Add sources if available
            if result.get('sources') and len(result['sources']) > 0:
                p = tf.add_paragraph()
                p.clear()
                run = p.add_run()
                run.text = "Sources:"
                p.alignment = PP_ALIGN.LEFT
                
                # Apply run formatting
                font = run.font
                font.bold = True
                font.name = 'Calibri'
                font.size = Pt(11)  # Match result font size
                font.color.rgb = THEME_PRIMARY
                
                for source in result['sources'][:3]:  # Limit to first 3 sources to save space
                    p = tf.add_paragraph()
                    p.clear()
                    run = p.add_run()
                    run.text = source[:100] + "..." if len(source) > 100 else source
                    
                    # Apply paragraph formatting
                    p.level = 1  # Indent sources
                    p.space_before = Pt(2)
                    p.space_after = Pt(2)
                    
                    # Apply run formatting
                    font = run.font
                    font.size = Pt(9)  # Smaller font for sources
                    font.name = 'Calibri'
                    font.italic = True
                    font.color.rgb = THEME_TEXT
        
        # Save the presentation to a BytesIO object
        pptx_buffer = BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        
        # Return the PowerPoint file as a downloadable file
        return StreamingResponse(
            pptx_buffer, 
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f"attachment; filename={company}_research.pptx"}
        )
    else:
        raise HTTPException(status_code=400, detail=f"Unsupported export format: {format}")